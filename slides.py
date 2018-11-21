# https://python-pptx.readthedocs.io/en/latest/user/presentations.html

import json
import requests
from pptx import Presentation
from cognition_api import get_slide_data, process_slide_data
import os.path

# title slide
# background
# motivation
# focus
# questions
# deliverables?
# data/methods/results (insert)
# value
# next steps

def move_to_front(prs):
  last_slide = len(prs.slides) - 1
  slides = list(prs.slides._sldIdLst)
  prs.slides._sldIdLst.remove(slides[last_slide])
  prs.slides._sldIdLst.insert(0, slides[last_slide])

def make_slides(prs):
  SLD_TITLE = 0
  SLD_HEAD_COPY = 1
  SLD_HEAD_BULLETS = 2
  # SLD_HEAD_SUBHEAD_COPY = 3
  # SLD_HEAD_ONLY = 7

  title_layout = prs.slide_layouts[SLD_TITLE]
  plain_layout = prs.slide_layouts[SLD_HEAD_COPY]
  bullet_layout = prs.slide_layouts[SLD_HEAD_BULLETS]

  ### deliverables

  slide = prs.slides.add_slide(bullet_layout)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[16]
  title_shape.text = 'Deliverables'
  tf = body_shape.text_frame
  items = d['deliverables']
  for item in items:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

  move_to_front(prs)

  ### questions

  slide = prs.slides.add_slide(bullet_layout)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[16]
  title_shape.text = 'Sprint Questions'
  tf = body_shape.text_frame
  items = d['sprint_question']
  for item in items:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

  move_to_front(prs)

  ### problem statement

  slide = prs.slides.add_slide(plain_layout)
  slide.placeholders[0].text = "Problem Statement"
  slide.placeholders[16].text = d['problem_statement']

  move_to_front(prs)

  ### motivation

  slide = prs.slides.add_slide(plain_layout)
  slide.placeholders[0].text = "Motivation"
  slide.placeholders[16].text = d['motivation']

  move_to_front(prs)

  ### background

  slide = prs.slides.add_slide(plain_layout)
  slide.placeholders[0].text = "Background"
  slide.placeholders[16].text = d['background']

  move_to_front(prs)

  ### title slide

  slide = prs.slides.add_slide(title_layout)
  slide.placeholders[0].text = 'Rally ' + d['sprint_id'] + ': ' + d['title']
  slide.placeholders[15].text = 'Completed ' + d['end_date'] # date
  slide.placeholders[17].text = 'Presented by ' + d['presenter'] + ' on behalf of rally participants ' + d['participants']

  move_to_front(prs)

  ### data/methods/results are already part of the deck

  ### key findings

  slide = prs.slides.add_slide(bullet_layout)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[16]
  title_shape.text = 'Key Findings'
  tf = body_shape.text_frame
  items = d['key_findings']
  for item in items:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

  ### value

  slide = prs.slides.add_slide(plain_layout)
  slide.placeholders[0].text = "Value"
  slide.placeholders[16].text = d['value']

  ### next steps

  slide = prs.slides.add_slide(bullet_layout)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[16]
  title_shape.text = 'Next Steps'
  tf = body_shape.text_frame
  items = d['next_steps']
  for item in items:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

  ### remove INSTRUCTIONS slide before saving

  slides = list(prs.slides)
  slides2 = list(prs.slides._sldIdLst)
  rm_idx = next((i for i in range(len(slides)) if slides[i].slide_layout.name == 'INSTRUCTIONS'), None)
  if rm_idx != None:
    prs.slides._sldIdLst.remove(slides2[rm_idx])

  return prs


# with open('mock.json') as f:
#   d = json.load(f)

with open('from_api.json') as f:
  a = json.load(f)
d = process_slide_data(a)

# d = get_slide_data('135')

pptx_path = 'template_ki_empty.pptx'
# if the slides URL is valid and we can download
# then we will use this slide deck, otherwise fall
# back on empty template
if d['ds_slides_url'] != '':
  url = d['ds_slides_url']
  r = requests.get(url)
  with open('/tmp/pres.pptx', 'wb') as f:  
    f.write(r.content)
  if r.status_code == 200 and os.path.isfile('/tmp/pres.pptx'):
    pptx_path = '/tmp/pres.pptx'

prs = Presentation(pptx_path)
# len(prs.slides)
# for layout in prs.slide_layouts:
#   print(layout.name)

prs = make_slides(prs)

prs.save('test.pptx')




# for shape in slide.placeholders:
#   print('%d %s' % (shape.placeholder_format.idx, shape.name))
# 0 Title 1
# 15 Text Placeholder 3
# 17 Text Placeholder 4

# for shape in slide.placeholders:
#   print('%d %s' % (shape.placeholder_format.idx, shape.name))
# # 0 Title 2
# # 16 Text Placeholder 1



# slide2 = prs.slides.add_slide(plain_layout)
# for shape in slide2.placeholders:
#   print('%d %s' % (shape.placeholder_format.idx, shape.name))
# # 0 Title 2
# # 16 Text Placeholder 1

# slide2.placeholders[0].text = "Test Title"
# body = slide2.placeholders[16]
# body.has_text_frame
# text_frame = body.text_frame
# text_frame.clear()
# text_frame.paragraphs[0]
# p = text_frame.paragraphs[0]
# run = p.add_run()
# run.text = 'Bold Text:'
# font = run.font
# font.bold = True
# # font.name = 'Calibri'
# # font.size = Pt(18)
# # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
# # run.hyperlink.address = 'https://github.com/scanny/python-pptx'
# p = text_frame.add_paragraph()
# p.text = "Text"
